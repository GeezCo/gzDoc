import io
from typing import Optional
from PyPDF2 import PdfReader
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation


class DocumentParser:
    """文档解析器"""

    @staticmethod
    def parse_pdf(file_data: bytes) -> str:
        """解析PDF文件"""
        try:
            pdf_file = io.BytesIO(file_data)
            reader = PdfReader(pdf_file)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            return text.strip()
        except Exception as e:
            raise Exception(f"PDF解析失败: {str(e)}")

    @staticmethod
    def parse_docx(file_data: bytes) -> str:
        """解析Word文档"""
        try:
            doc_file = io.BytesIO(file_data)
            doc = Document(doc_file)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            return text.strip()
        except Exception as e:
            raise Exception(f"Word文档解析失败: {str(e)}")

    @staticmethod
    def parse_xlsx(file_data: bytes) -> str:
        """解析Excel文件"""
        try:
            excel_file = io.BytesIO(file_data)
            workbook = load_workbook(excel_file, read_only=True)
            text = ""
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"\n=== {sheet_name} ===\n"
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                    text += row_text + "\n"
            return text.strip()
        except Exception as e:
            raise Exception(f"Excel文件解析失败: {str(e)}")

    @staticmethod
    def parse_pptx(file_data: bytes) -> str:
        """解析PowerPoint文件"""
        try:
            ppt_file = io.BytesIO(file_data)
            prs = Presentation(ppt_file)
            text = ""
            for i, slide in enumerate(prs.slides, 1):
                text += f"\n=== Slide {i} ===\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip()
        except Exception as e:
            raise Exception(f"PowerPoint文件解析失败: {str(e)}")

    @staticmethod
    def parse_txt(file_data: bytes) -> str:
        """解析文本文件"""
        try:
            return file_data.decode("utf-8")
        except UnicodeDecodeError:
            try:
                return file_data.decode("gbk")
            except Exception as e:
                raise Exception(f"文本文件解析失败: {str(e)}")

    @classmethod
    def parse(cls, file_data: bytes, file_type: str) -> str:
        """根据文件类型解析文档"""
        file_type = file_type.lower()

        parsers = {
            "pdf": cls.parse_pdf,
            "docx": cls.parse_docx,
            "doc": cls.parse_docx,
            "xlsx": cls.parse_xlsx,
            "xls": cls.parse_xlsx,
            "pptx": cls.parse_pptx,
            "ppt": cls.parse_pptx,
            "txt": cls.parse_txt,
        }

        parser = parsers.get(file_type)
        if not parser:
            raise Exception(f"不支持的文件类型: {file_type}")

        return parser(file_data)
